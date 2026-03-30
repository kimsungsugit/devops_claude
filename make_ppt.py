"""
AI 활용 아키텍처 PPT - White Modern Design + 설명 텍스트
실행: C:/msys64/mingw64/bin/python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── 색상 팔레트 ──────────────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BG         = RGBColor(0xF4, 0xF6, 0xF9)
C_CARD       = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
C_TEXT_MID   = RGBColor(0x44, 0x44, 0x66)
C_TEXT_LIGHT = RGBColor(0x88, 0x88, 0xAA)

A1 = RGBColor(0x41, 0x6F, 0xE0)   # 인디고블루
A2 = RGBColor(0x0F, 0xB5, 0x8A)   # 에메랄드
A3 = RGBColor(0x7C, 0x5C, 0xF6)   # 바이올렛
A4 = RGBColor(0xF5, 0x6E, 0x00)   # 앰버오렌지
A5 = RGBColor(0xE0, 0x31, 0x7E)   # 핑크

C_TITLE_BG  = RGBColor(0x1A, 0x1A, 0x2E)
C_TITLE_SUB = RGBColor(0x41, 0x6F, 0xE0)
C_SUMMARY   = RGBColor(0xEC, 0xEF, 0xF8)   # 요약 스트립 배경

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_rr(slide, x, y, w, h, fill, line_color=None, line_w=Pt(0.75), adj=0.04):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    shape.adjustments[0] = adj
    return shape


def add_tb(slide, x, y, w, h, text, size=Pt(10), bold=False,
           color=C_TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return txBox


def add_tb2(slide, x, y, w, h, line1, line2,
            size1=Pt(10), size2=Pt(8.5),
            bold1=True, bold2=False,
            color1=C_WHITE, color2=C_WHITE,
            align=PP_ALIGN.CENTER, italic2=False):
    """두 줄 텍스트박스"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = align
    r1 = p1.add_run()
    r1.text = line1
    r1.font.size = size1
    r1.font.bold = bold1
    r1.font.color.rgb = color1
    r1.font.name = "Malgun Gothic"

    p2 = tf.add_paragraph()
    p2.alignment = align
    r2 = p2.add_run()
    r2.text = line2
    r2.font.size = size2
    r2.font.bold = bold2
    r2.font.italic = italic2
    r2.font.color.rgb = color2
    r2.font.name = "Malgun Gothic"

    return txBox


def add_arrow(slide, x1, y1, x2, y2, color=C_TEXT_LIGHT, width=Pt(1.5)):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


# ─────────────────────────────────────────────────────────────
def make_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 전체 배경 ─────────────────────────────────────────────
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)

    # ═══════════════════════════════════════════════════════════
    #  타이틀 영역 (2줄 구조로 확장 — 0 ~ 0.95")
    # ═══════════════════════════════════════════════════════════
    TITLE_H = Inches(0.95)
    add_rect(slide, 0, 0, SLIDE_W, TITLE_H, C_TITLE_BG)
    add_rect(slide, 0, 0, Inches(0.07), TITLE_H, C_TITLE_SUB)   # 왼쪽 강조선

    # 메인 타이틀
    add_tb(slide, Inches(0.25), Inches(0.06), Inches(8.5), Inches(0.46),
           "AI 기반 DevOps 자동화 아키텍처",
           size=Pt(21), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # 서브 타이틀 (흐름 태그)
    add_tb(slide, Inches(0.25), Inches(0.52), Inches(10), Inches(0.36),
           "소스코드 변경 감지  →  RAG 지식베이스 구성  →  LLM 문서 추론  →  설계·시험 문서 자동 생성  →  AI 어시스턴트 조율",
           size=Pt(9.5), color=rgb(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT)

    # 우측 키워드 뱃지
    add_tb(slide, Inches(10.2), Inches(0.52), Inches(2.9), Inches(0.34),
           "LLM · RAG · VectorDB · Auto-Doc",
           size=Pt(9), color=rgb(0x88, 0xAA, 0xFF), align=PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════════════════════════
    #  전체 요약 스트립 (0.95" ~ 1.42")
    # ═══════════════════════════════════════════════════════════
    SUMMARY_Y = TITLE_H
    SUMMARY_H = Inches(0.47)
    add_rect(slide, 0, SUMMARY_Y, SLIDE_W, SUMMARY_H, C_SUMMARY)
    # 좌측 강조 바
    add_rect(slide, 0, SUMMARY_Y, Inches(0.07), SUMMARY_H, C_TITLE_SUB)
    # 구분선 (하단)
    add_rect(slide, 0, SUMMARY_Y + SUMMARY_H - Inches(0.018),
             SLIDE_W, Inches(0.018), rgb(0xD0, 0xD5, 0xE8))

    add_tb(slide, Inches(0.22), SUMMARY_Y + Inches(0.06),
           Inches(1.2), Inches(0.32),
           "시스템 개요",
           size=Pt(9), bold=True, color=C_TITLE_SUB, align=PP_ALIGN.LEFT)

    add_tb(slide, Inches(1.42), SUMMARY_Y + Inches(0.06),
           Inches(11.7), Inches(0.32),
           "Git/SVN 변경을 감지하면, 관련 문서를 벡터 DB에서 검색(RAG)하여 LLM 컨텍스트로 주입하고, "
           "변경 유형에 맞는 설계서(UDS)·시험 명세(STS/SUTS/SITS)를 자동 생성합니다. "
           "LangGraph 기반 AI 어시스턴트가 전체 워크플로우를 조율하며 근거(Evidence)와 함께 결과를 제공합니다.",
           size=Pt(9), color=C_TEXT_MID, align=PP_ALIGN.LEFT, wrap=True)

    # ═══════════════════════════════════════════════════════════
    #  열 레이아웃 정의
    # ═══════════════════════════════════════════════════════════
    TOP   = SUMMARY_Y + SUMMARY_H + Inches(0.1)
    ROW_H = Inches(5.48)
    COL_W = Inches(2.42)
    GAP   = Inches(0.165)
    COLS_X = [Inches(0.15) + (COL_W + GAP) * i for i in range(5)]

    accents = [A1, A2, A3, A4, A5]

    col_info = [
        ("① 소스 변경 감지",           "변경 감지 → 분류 → 영향도"),
        ("② RAG / Knowledge Base",     "문서 수집 → 임베딩 → 검색"),
        ("③ LLM 엔진",                 "추론 · 생성 · 프롬프트 관리"),
        ("④ 문서 자동 생성",           "UDS · STS · SUTS · SITS"),
        ("⑤ 워크플로우 오케스트레이션", "LangGraph · MCP · 채팅 AI"),
    ]

    HEADER_H = Inches(0.65)

    for i, cx in enumerate(COLS_X):
        ac = accents[i]

        # 열 그림자
        add_rr(slide, cx + Inches(0.03), TOP + Inches(0.03),
               COL_W, ROW_H, rgb(0xC0, 0xC8, 0xDC), adj=0.03)

        # 열 배경 (흰색)
        add_rr(slide, cx, TOP, COL_W, ROW_H, C_WHITE,
               line_color=rgb(0xDD, 0xDD, 0xEE), line_w=Pt(0.5), adj=0.03)

        # 열 헤더 (Accent + 부제)
        add_rr(slide, cx, TOP, COL_W, HEADER_H, ac, adj=0.03)
        add_rect(slide, cx, TOP + Inches(0.32), COL_W, HEADER_H - Inches(0.32), ac)

        title_txt, sub_txt = col_info[i]

        # 열 제목 (라인 1)
        add_tb(slide, cx + Inches(0.08), TOP + Inches(0.04),
               COL_W - Inches(0.16), Inches(0.3),
               title_txt, size=Pt(10), bold=True,
               color=C_WHITE, align=PP_ALIGN.CENTER)

        # 열 부제 (라인 2) — 연한 흰색 italic
        add_tb(slide, cx + Inches(0.08), TOP + Inches(0.34),
               COL_W - Inches(0.16), Inches(0.25),
               sub_txt, size=Pt(8), bold=False, italic=True,
               color=rgb(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

    # ── 화살표 (열 사이) ──────────────────────────────────────
    arrow_y = TOP + ROW_H * 0.52
    for i in range(4):
        ax_s = COLS_X[i] + COL_W + Inches(0.02)
        ax_e = COLS_X[i + 1] - Inches(0.02)
        add_arrow(slide, ax_s, arrow_y, ax_e, arrow_y,
                  color=accents[i], width=Pt(2.0))
        add_tb(slide, ax_e - Inches(0.03), arrow_y - Inches(0.12),
               Inches(0.12), Inches(0.24), "▶",
               size=Pt(8), color=accents[i], align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════
    #  각 열 카드 내용
    # ═══════════════════════════════════════════════════════════
    CARD_START = TOP + HEADER_H + Inches(0.06)

    _draw_cards(slide, COLS_X[0], CARD_START, COL_W, A1,
                ROW_H - HEADER_H - Inches(0.1), [
        ("Git / SVN",
         "브랜치 커밋 감지\n변경 파일 목록 추출\nSVN Working Copy 지원"),
        ("SCM Registry",
         "프로젝트 메타 관리\nscm_registry.json\n문서 경로 자동 매핑"),
        ("변경 분류",
         "SIGNATURE · BODY\nNEW · DELETE\nVARIABLE · HEADER"),
        ("Impact 분석",
         "함수 호출 트리 분석\n최대 2-hop 영향 탐색\n모듈 신뢰도 스코어링"),
    ])

    _draw_cards(slide, COLS_X[1], CARD_START, COL_W, A2,
                ROW_H - HEADER_H - Inches(0.1), [
        ("문서 수집 · 청킹",
         "DOCX (헤딩 기준)\nC/C++ (AST 함수 단위)\nTXT 1,200자 + 200자 오버랩"),
        ("임베딩",
         "외부 HTTP API\nSentence-Transformers\nLRU 캐시 1,000건"),
        ("Vector DB",
         "SQLite (기본)\nPostgreSQL + pgvector\n64-dim 벡터 저장"),
        ("유사도 검색",
         "코사인 유사도\n프로젝트 · 최근성 · 태그\n다중 부스팅 전략"),
    ])

    _draw_cards(slide, COLS_X[2], CARD_START, COL_W, A3,
                ROW_H - HEADER_H - Inches(0.1), [
        ("멀티 프로바이더",
         "Google Gemini  ★ 주력\nOpenAI GPT\nAnthropic Claude\nOllama (로컬)"),
        ("LLM Adapter",
         "추상화 어댑터 패턴\nTemperature 0.3\nMax Token 8,192\nContext Caching 지원"),
        ("Prompt 관리",
         "분석 · 작성 · 검토\n감사 · 로직 · 섹션\n→ 6개 전용 프롬프트 파일"),
    ])

    _draw_cards(slide, COLS_X[3], CARD_START, COL_W, A4,
                ROW_H - HEADER_H - Inches(0.1), [
        ("UDS",
         "단위 설계 명세서\n개요 · 인터페이스 · 로직\nEvidence 근거 추적"),
        ("STS",
         "소프트웨어 시험 명세\n배치 병렬 처리\n600+ 함수 자동화"),
        ("SUTS / SITS",
         "단위 · 통합 시험 명세\n변경 영향도 기반\n자동 대상 결정"),
        ("생성 매트릭스",
         "변경유형 × 문서종류\nAUTO / FLAG 자동 결정\n다단계 Retry 포함"),
    ])

    _draw_cards(slide, COLS_X[4], CARD_START, COL_W, A5,
                ROW_H - HEADER_H - Inches(0.1), [
        ("LangGraph",
         "상태 머신 기반\n의도 분석 → Tool 선택\n실행 → 근거 → 응답"),
        ("MCP Bridge",
         "코드 검색 / Git 조작\nJenkins 빌드 분석\n문서 검색 통합"),
        ("채팅 어시스턴트",
         "RAG 기반 Q&A\n근거 · 인용 자동 추적\n승인 워크플로우"),
    ])

    # ═══════════════════════════════════════════════════════════
    #  하단 범례 바
    # ═══════════════════════════════════════════════════════════
    LEG_Y = TOP + ROW_H + Inches(0.1)
    add_rect(slide, 0, LEG_Y, SLIDE_W, Inches(0.4), rgb(0xE2, 0xE6, 0xF3))
    add_rect(slide, 0, LEG_Y, Inches(0.07), Inches(0.4), C_TITLE_SUB)

    legend_items = [
        (A1, "① 소스 변경 감지"),
        (A2, "② RAG / Vector DB"),
        (A3, "③ LLM 엔진"),
        (A4, "④ 문서 자동 생성"),
        (A5, "⑤ 워크플로우"),
    ]
    lx = Inches(0.35)
    for lc, lt in legend_items:
        add_rr(slide, lx, LEG_Y + Inches(0.1), Inches(0.2), Inches(0.2), lc, adj=0.5)
        add_tb(slide, lx + Inches(0.26), LEG_Y + Inches(0.08),
               Inches(2.0), Inches(0.26), lt,
               size=Pt(9), color=C_TEXT_MID)
        lx += Inches(2.3)

    return slide


def _draw_cards(slide, col_x, start_y, col_w, accent, available_h, items):
    n = len(items)
    card_h = available_h / n - Inches(0.065)
    pad = Inches(0.1)
    iw = col_w - pad * 2

    for i, (title, body) in enumerate(items):
        iy = start_y + i * (card_h + Inches(0.065))
        ix = col_x + pad

        # 카드 그림자
        add_rr(slide, ix + Inches(0.022), iy + Inches(0.022),
               iw, card_h, rgb(0xC8, 0xCC, 0xDC), adj=0.04)

        # 카드 본체
        add_rr(slide, ix, iy, iw, card_h, C_WHITE,
               line_color=rgb(0xE0, 0xE4, 0xF0), line_w=Pt(0.5), adj=0.04)

        # 좌측 Accent 세로 바
        add_rect(slide, ix, iy + Inches(0.07),
                 Inches(0.05), card_h - Inches(0.14), accent)

        # 카드 제목
        add_tb(slide, ix + Inches(0.1), iy + Inches(0.05),
               iw - Inches(0.12), Inches(0.29),
               title, size=Pt(9.5), bold=True,
               color=accent, align=PP_ALIGN.LEFT)

        # 구분선
        add_rect(slide, ix + Inches(0.1), iy + Inches(0.32),
                 iw - Inches(0.18), Inches(0.01),
                 rgb(0xE8, 0xEA, 0xF4))

        # 카드 본문
        add_tb(slide, ix + Inches(0.1), iy + Inches(0.34),
               iw - Inches(0.13), card_h - Inches(0.4),
               body, size=Pt(8.5), color=C_TEXT_MID,
               align=PP_ALIGN.LEFT, wrap=True)


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    make_slide(prs)

    out_path = "d:/Project/devops/260105/AI_Architecture.pptx"
    prs.save(out_path)
    print(f"저장 완료: {out_path}")
